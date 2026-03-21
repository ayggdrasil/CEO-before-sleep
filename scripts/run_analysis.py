#!/usr/bin/env python3
"""
CEO Dream Works - Multi-Agent Business Analysis Execution Script

This script reads all materials in a target folder,
performs sequential analysis across 13 dimensions × 5 agents,
and generates a final markdown report.

Usage:
    python scripts/run_analysis.py --folder /path/to/project/ --output /path/to/output/CEO_DreamWorks_Report.md
"""

import os
import sys
import json
import argparse
from datetime import datetime
from pathlib import Path

# ── Analysis Dimensions ─────────────────────────────────────
ANALYSIS_ITEMS = [
    {"id": 1,  "name": "Item Profitability",    "desc": "Business model, unit economics, margin structure, scalability"},
    {"id": 2,  "name": "Talent Acquisition",    "desc": "Key talent pipeline, hiring strategy, retention, compensation structure"},
    {"id": 3,  "name": "Team Composition",      "desc": "Current team capabilities, gap analysis, org structure, decision-making framework"},
    {"id": 4,  "name": "Technology",            "desc": "Tech stack, technical moat/differentiation, tech debt, scalability"},
    {"id": 5,  "name": "Legal",                 "desc": "Regulatory risk, IP protection, contract structure, compliance"},
    {"id": 6,  "name": "Corporate Structure",   "desc": "Entity type, equity structure, governance, global structure"},
    {"id": 7,  "name": "Finance",               "desc": "Financial health, cash flow, burn rate, funding history/plans"},
    {"id": 8,  "name": "Execution",             "desc": "Roadmap, milestones, execution capability, operational efficiency"},
    {"id": 9,  "name": "GTM",                   "desc": "Market entry strategy, channels, partnerships, early customers"},
    {"id": 10, "name": "Sales",                 "desc": "Sales pipeline, conversion rates, sales cycle, pricing"},
    {"id": 11, "name": "Most Urgent Task",      "desc": "The #1 thing that must be solved right now"},
    {"id": 12, "name": "Biggest Risk",          "desc": "The #1 thing that can kill this business"},
    {"id": 13, "name": "Crazy Founder's Advice","desc": "Path to $100B from a serial founder with multiple exits"},
    {"id": 14, "name": "Execution Retro",       "desc": "Planned vs. shipped velocity over the last 90 days (post-seed only)"},
]

AGENTS = [
    {"id": "reviewer",          "name": "Reviewer",          "emoji": "📋", "desc": "Fact-based status assessment"},
    {"id": "problem_finder",    "name": "Problem Finder",    "emoji": "⚠️", "desc": "Hidden issues, risks, and gaps"},
    {"id": "problem_solver",    "name": "Problem Solver",    "emoji": "💡", "desc": "Concrete solutions and alternatives"},
    {"id": "researcher",        "name": "Researcher",        "emoji": "📊", "desc": "Market benchmark comparison"},
    {"id": "consultant",        "name": "Consultant",        "emoji": "🎯", "desc": "Final CEO-level advice"},
    {"id": "competitor_shadow", "name": "Competitor Shadow", "emoji": "🕵️", "desc": "How competitors exploit this weakness"},
]

# ── File Reading Utilities ──────────────────────────────────
READABLE_EXTENSIONS = {
    '.md', '.txt', '.csv', '.json', '.yaml', '.yml', '.toml',
    '.py', '.js', '.ts', '.html', '.css', '.sh',
    '.log', '.xml', '.ini', '.cfg', '.conf',
}

BINARY_EXTENSIONS = {
    '.pdf': 'PDF',
    '.docx': 'Word',
    '.xlsx': 'Excel',
    '.pptx': 'PowerPoint',
}


def read_text_file(filepath):
    """Read a plain text file."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        return content[:50000]  # Max 50K chars per file
    except Exception as e:
        return f"[Read failed: {e}]"


def read_pdf(filepath):
    """Extract text from PDF using PyMuPDF or pdftotext fallback."""
    try:
        import subprocess
        result = subprocess.run(
            ['python3', '-c', f'''
import fitz
doc = fitz.open("{filepath}")
for page in doc:
    print(page.get_text())
'''],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout[:50000]

        # pdftotext fallback
        result = subprocess.run(
            ['pdftotext', filepath, '-'],
            capture_output=True, text=True, timeout=30
        )
        return result.stdout[:50000] if result.returncode == 0 else "[PDF text extraction failed]"
    except Exception as e:
        return f"[PDF read failed: {e}]"


def read_docx(filepath):
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(filepath)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return text[:50000]
    except Exception as e:
        return f"[DOCX read failed: {e}]"


def read_xlsx(filepath):
    """Extract content from XLSX."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, data_only=True)
        output = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            output.append(f"\n--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                output.append(' | '.join(cells))
        return '\n'.join(output)[:50000]
    except Exception as e:
        return f"[XLSX read failed: {e}]"


def read_pptx(filepath):
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return '\n'.join(text_parts)[:50000]
    except Exception as e:
        return f"[PPTX read failed: {e}]"


def collect_files(folder_path):
    """Recursively collect all files in folder and extract their content."""
    context_parts = []
    file_list = []

    for root, dirs, files in os.walk(folder_path):
        # Skip hidden directories
        dirs[:] = [d for d in dirs if not d.startswith('.')]

        for fname in sorted(files):
            if fname.startswith('.'):
                continue

            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, folder_path)
            ext = os.path.splitext(fname)[1].lower()

            file_list.append(rel_path)

            content = None
            if ext in READABLE_EXTENSIONS:
                content = read_text_file(filepath)
            elif ext == '.pdf':
                content = read_pdf(filepath)
            elif ext == '.docx':
                content = read_docx(filepath)
            elif ext == '.xlsx':
                content = read_xlsx(filepath)
            elif ext == '.pptx':
                content = read_pptx(filepath)
            else:
                content = f"[Unsupported file type: {ext}]"

            if content:
                context_parts.append(f"\n{'='*60}\nFile: {rel_path}\n{'='*60}\n{content}\n")

    return '\n'.join(context_parts), file_list


def generate_report_skeleton(project_name, file_list, timestamp):
    """Generate the final markdown report skeleton."""

    report = f"""# CEO Dream Works Analysis Report

> **Target**: {project_name}
> **Date**: {timestamp}
> **Model**: claude-opus-4-6
> **Analysis loops**: {len(ANALYSIS_ITEMS)} dimensions × {len(AGENTS)} agents = {len(ANALYSIS_ITEMS) * len(AGENTS)} analyses
> **Files analyzed**: {len(file_list)}

---

## Executive Summary

[To be completed after analysis]

## Overall Grade

[To be completed after analysis]

---

"""

    for item in ANALYSIS_ITEMS:
        report += f"## {item['id']}. {item['name']}\n\n"
        report += f"*{item['desc']}*\n\n"

        for agent in AGENTS:
            report += f"### {agent['emoji']} {agent['desc']} ({agent['name']})\n\n"
            report += "[Awaiting analysis]\n\n"

        report += "---\n\n"

    # Appendix
    report += """## Appendix

### A. List of Materials Analyzed

"""
    for f in file_list:
        report += f"- `{f}`\n"

    report += """
### B. Methodology

This report was generated using the CEO Dream Works skill.
Five expert sub-agents (Reviewer, Problem Finder, Problem Solver, Researcher, Consultant)
sequentially analyzed 13 key dimensions. Each agent referenced prior agents' results
to cumulatively deepen the analysis.

### C. Disclaimer

This report is a reference document generated by an AI analysis tool.
For legal, financial, and tax matters, please consult qualified professionals.
"""

    return report


def main():
    parser = argparse.ArgumentParser(description='CEO Dream Works Analysis — Phase 0 Runner')
    parser.add_argument('--folder', required=True, help='Path to the target folder')
    parser.add_argument('--output', default='CEO_DreamWorks_Report.md', help='Path for the output report')
    parser.add_argument('--project-name', default=None, help='Project name (defaults to folder name)')

    args = parser.parse_args()

    folder = args.folder
    output = args.output
    project_name = args.project_name or os.path.basename(os.path.normpath(folder))
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    print(f"Target folder : {folder}")
    print(f"Output path   : {output}")
    print(f"Project name  : {project_name}")

    # Phase 1: Material collection (Phase 0 = Scope Alignment questions, asked by Claude before this script)
    print("\nPhase 1: Collecting materials...")
    context, file_list = collect_files(folder)
    print(f"  {len(file_list)} files found")
    print(f"  Context size : {len(context):,} characters")

    # Save context bundle
    output_dir = os.path.dirname(os.path.abspath(output))
    os.makedirs(output_dir, exist_ok=True)

    bundle_path = os.path.join(output_dir, 'context_bundle.md')
    with open(bundle_path, 'w', encoding='utf-8') as f:
        f.write(context)
    print(f"  Context bundle saved: {bundle_path}")

    # Generate report skeleton
    report = generate_report_skeleton(project_name, file_list, timestamp)
    with open(output, 'w', encoding='utf-8') as f:
        f.write(report)
    print(f"\nReport skeleton saved: {output}")

    # Save metadata
    meta = {
        "project_name": project_name,
        "timestamp": timestamp,
        "folder": folder,
        "output": output,
        "file_count": len(file_list),
        "file_list": file_list,
        "context_size_chars": len(context),
        "items": ANALYSIS_ITEMS,
        "agents": AGENTS,
        "total_loops": len(ANALYSIS_ITEMS) * len(AGENTS),
    }
    meta_path = os.path.join(output_dir, 'analysis_meta.json')
    with open(meta_path, 'w', encoding='utf-8') as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    print(f"Metadata saved: {meta_path}")

    print(f"\n{'='*60}")
    print(f"Phase 1 complete. Claude will now begin Phase 2 (sequential agent analysis).")
    print(f"  Note: Run Phase 0 (Scope Alignment) first — ask the 3 forcing questions.")
    print(f"Dimensions : {len(ANALYSIS_ITEMS)}")
    print(f"Agents     : {len(AGENTS)}")
    print(f"Total loops: {len(ANALYSIS_ITEMS) * len(AGENTS)}")
    print(f"{'='*60}")


if __name__ == '__main__':
    main()
